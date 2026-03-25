# coding=utf-8
from pptx import Presentation
import language_tool_python
import difflib

# Load tools
tool = language_tool_python.LanguageTool('en-US')
prs = Presentation("C:/Users/jmayhall/Downloads/Mayhall_AMS_Tropical.pptx")


def correct_text(text):
    matches = tool.check(text)
    return language_tool_python.utils.correct(text, matches), matches


def show_diff(original, corrected):
    diff = difflib.ndiff(original.split(), corrected.split())

    added = []
    removed = []

    for token in diff:
        if token.startswith("+ "):
            added.append(token[2:])
        elif token.startswith("- "):
            removed.append(token[2:])

    return added, removed, diff


for slide_num, slide in enumerate(prs.slides, start=1):
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para_num, paragraph in enumerate(shape.text_frame.paragraphs, start=1):
                original = paragraph.text.strip()
                if not original:
                    continue

                corrected, matches = correct_text(original)

                if matches:
                    added, removed, diff = show_diff(original, corrected)

                    print(f"\n--- Slide {slide_num}, Paragraph {para_num} ---")
                    print("Original :", original)
                    print("Corrected:", corrected)

                    print("\nChanges:")
                    print("Removed:", removed)
                    print("Added  :", added)

                    print("\nFull Diff:")
                    print(" ".join(diff))