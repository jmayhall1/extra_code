# coding=utf-8
from pptx import Presentation
import difflib

def extract_slide_text(prs):
    """Return a list where each item is all text from one slide."""
    slides_text = []

    for slide in prs.slides:
        slide_text = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        slide_text.append(text)
        slides_text.append("\n".join(slide_text))

    return slides_text


def compare_presentations(file1, file2):
    prs1 = Presentation(file1)
    prs2 = Presentation(file2)

    slides1 = extract_slide_text(prs1)
    slides2 = extract_slide_text(prs2)

    max_len = max(len(slides1), len(slides2))
    different_slides = []

    for i in range(max_len):
        text1 = slides1[i] if i < len(slides1) else ""
        text2 = slides2[i] if i < len(slides2) else ""

        if text1 != text2:
            different_slides.append(i + 1)  # Slide numbers start at 1

    return different_slides


# === USAGE ===
file1 = "C:/Users/jmayhall/Downloads/Module13-Agriculture-Urbanization-AES103-SP26.pptx"
file2 = "C:/Users/jmayhall/Downloads/Module13-Agriculture-Urbanization-AES103-SP26 (1).pptx"

diff_slides = compare_presentations(file1, file2)

print("Slides with differences:")
print(diff_slides)